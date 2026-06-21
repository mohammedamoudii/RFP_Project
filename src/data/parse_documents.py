"""Parse manifest-listed source files into normalized text elements.

This module reads the file manifest produced by ``create_manifest.py`` and
extracts text from supported document formats. Each parsed element keeps the
manifest metadata needed by later cleaning, chunking, retrieval, and evidence
validation stages.
"""
from __future__ import annotations

from pathlib import Path
import argparse
import json
import traceback
from typing import Any, Iterable

import pandas as pd
from tqdm import tqdm

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


# -----------------------------------------------------------------------------
# Project paths and parser configuration
# -----------------------------------------------------------------------------
MANIFEST_PATH = Path("data/processed/file_manifest.csv")
OUTPUT_JSONL = Path("data/processed/parsed_elements.jsonl")
ERRORS_CSV = Path("data/processed/parsing_errors.csv")

SUPPORTED_FILE_TYPES = {
    "pdf",
    "docx",
    "pptx",
    "md",
    "txt",
    "xlsx",
}


# -----------------------------------------------------------------------------
# Text and JSONL helpers
# -----------------------------------------------------------------------------
def safe_text(value: Any) -> str:
    """Convert nullable values from files or DataFrames into clean strings."""
    if value is None:
        return ""

    try:
        if pd.isna(value):
            return ""
    except (TypeError, ValueError):
        pass

    return str(value).strip()


def optional_text(value: Any) -> str | None:
    """Return cleaned text, using None for empty optional metadata fields."""
    text = safe_text(value)
    return text or None


def write_jsonl(
    rows: Iterable[dict[str, Any]],
    output_path: Path,
) -> None:
    """Write dictionaries as UTF-8 JSONL records."""
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with output_path.open("w", encoding="utf-8") as file:
        for row in rows:
            file.write(
                json.dumps(row, ensure_ascii=False) + "\n"
            )


def read_jsonl(path: Path) -> list[dict[str, Any]]:
    """Read a JSONL file into dictionaries, skipping empty lines."""
    if not path.exists():
        return []

    rows: list[dict[str, Any]] = []

    with path.open("r", encoding="utf-8") as file:
        for line_number, line in enumerate(file, start=1):
            text = line.strip()

            if not text:
                continue

            try:
                value = json.loads(text)
            except json.JSONDecodeError as error:
                raise ValueError(
                    f"Invalid JSONL in {path} at line "
                    f"{line_number}: {error}"
                ) from error

            if isinstance(value, dict):
                rows.append(value)

    return rows


# -----------------------------------------------------------------------------
# Metadata and document parsers
# -----------------------------------------------------------------------------
def base_metadata(row: dict[str, Any] | pd.Series) -> dict[str, Any]:
    """Extract shared manifest metadata for every parsed element."""
    return {
        "file_id": safe_text(row.get("file_id")),
        "database_target": safe_text(
            row.get("database_target")
        ),
        "collection_name": safe_text(
            row.get("collection_name")
        ),
        "opportunity_id": optional_text(
            row.get("opportunity_id")
        ),
        "project_id": optional_text(
            row.get("project_id")
        ),
        "project_name": optional_text(
            row.get("project_name")
        ),
        "folder_name": safe_text(
            row.get("folder_name")
        ),
        "source_file": safe_text(
            row.get("source_file")
        ),
        "file_type": safe_text(
            row.get("file_type")
        ).lower(),
        "relative_path": safe_text(
            row.get("relative_path")
        ),
        "absolute_path": safe_text(
            row.get("absolute_path")
        ),
    }


def parse_pdf(
    file_path: Path,
    row: dict[str, Any] | pd.Series,
) -> list[dict[str, Any]]:
    """Extract one text element per readable PDF page."""
    elements: list[dict[str, Any]] = []
    metadata = base_metadata(row)

    reader = PdfReader(str(file_path))

    for page_index, page in enumerate(
        reader.pages,
        start=1,
    ):
        text = safe_text(page.extract_text())

        if not text:
            continue

        elements.append(
            {
                **metadata,
                "element_id": (
                    f"{metadata['file_id']}_page_{page_index}"
                ),
                "element_type": "pdf_page",
                "page_number": page_index,
                "slide_number": None,
                "sheet_name": None,
                "row_start": None,
                "row_end": None,
                "section": None,
                "content": text,
            }
        )

    return elements


def parse_docx(
    file_path: Path,
    row: dict[str, Any] | pd.Series,
) -> list[dict[str, Any]]:
    """Extract paragraph and table text from a Word document."""
    elements: list[dict[str, Any]] = []
    metadata = base_metadata(row)

    document = Document(str(file_path))
    element_index = 0

    for paragraph in document.paragraphs:
        text = safe_text(paragraph.text)

        if not text:
            continue

        element_index += 1

        elements.append(
            {
                **metadata,
                "element_id": (
                    f"{metadata['file_id']}_paragraph_"
                    f"{element_index}"
                ),
                "element_type": "docx_paragraph",
                "page_number": None,
                "slide_number": None,
                "sheet_name": None,
                "row_start": None,
                "row_end": None,
                "section": None,
                "content": text,
            }
        )

    for table_index, table in enumerate(
        document.tables,
        start=1,
    ):
        rows_text: list[str] = []

        for table_row in table.rows:
            cells = [
                safe_text(cell.text).replace("\n", " ")
                for cell in table_row.cells
            ]

            if any(cells):
                rows_text.append(" | ".join(cells))

        if rows_text:
            elements.append(
                {
                    **metadata,
                    "element_id": (
                        f"{metadata['file_id']}_table_"
                        f"{table_index}"
                    ),
                    "element_type": "docx_table",
                    "page_number": None,
                    "slide_number": None,
                    "sheet_name": None,
                    "row_start": None,
                    "row_end": None,
                    "section": None,
                    "content": "\n".join(rows_text),
                }
            )

    return elements


def parse_pptx(
    file_path: Path,
    row: dict[str, Any] | pd.Series,
) -> list[dict[str, Any]]:
    """Extract text and table content from each PowerPoint slide."""
    elements: list[dict[str, Any]] = []
    metadata = base_metadata(row)

    presentation = Presentation(str(file_path))

    for slide_index, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_text_parts: list[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = safe_text(shape.text)

                if text:
                    slide_text_parts.append(text)

            if getattr(shape, "has_table", False):
                table_rows: list[str] = []

                for table_row in shape.table.rows:
                    cells = [
                        safe_text(cell.text).replace("\n", " ")
                        for cell in table_row.cells
                    ]

                    if any(cells):
                        table_rows.append(" | ".join(cells))

                if table_rows:
                    slide_text_parts.append(
                        "\n".join(table_rows)
                    )

        slide_text = "\n\n".join(
            slide_text_parts
        ).strip()

        if slide_text:
            elements.append(
                {
                    **metadata,
                    "element_id": (
                        f"{metadata['file_id']}_slide_"
                        f"{slide_index}"
                    ),
                    "element_type": "pptx_slide",
                    "page_number": None,
                    "slide_number": slide_index,
                    "sheet_name": None,
                    "row_start": None,
                    "row_end": None,
                    "section": None,
                    "content": slide_text,
                }
            )

    return elements


def parse_text_file(
    file_path: Path,
    row: dict[str, Any] | pd.Series,
) -> list[dict[str, Any]]:
    """Read a Markdown or plain-text file as one parsed element."""
    metadata = base_metadata(row)

    text = file_path.read_text(
        encoding="utf-8",
        errors="ignore",
    ).strip()

    if not text:
        return []

    return [
        {
            **metadata,
            "element_id": f"{metadata['file_id']}_text_1",
            "element_type": "text_file",
            "page_number": None,
            "slide_number": None,
            "sheet_name": None,
            "row_start": None,
            "row_end": None,
            "section": None,
            "content": text,
        }
    ]


def dataframe_to_text(dataframe: pd.DataFrame) -> str:
    """Convert a spreadsheet slice into line-based pipe-delimited text."""
    dataframe = (
        dataframe
        .dropna(how="all")
        .dropna(axis=1, how="all")
    )

    if dataframe.empty:
        return ""

    dataframe = dataframe.fillna("")
    lines: list[str] = []

    for _, row in dataframe.iterrows():
        values = [
            safe_text(value).replace("\n", " ")
            for value in row.tolist()
        ]

        if any(values):
            lines.append(" | ".join(values))

    return "\n".join(lines).strip()


def parse_xlsx(
    file_path: Path,
    row: dict[str, Any] | pd.Series,
    rows_per_element: int = 40,
) -> list[dict[str, Any]]:
    """Extract spreadsheet sheets into row-window text elements."""
    elements: list[dict[str, Any]] = []
    metadata = base_metadata(row)

    excel_file = pd.ExcelFile(file_path)

    for sheet_name in excel_file.sheet_names:
        dataframe = pd.read_excel(
            excel_file,
            sheet_name=sheet_name,
            header=None,
            dtype=str,
        )

        dataframe = (
            dataframe
            .dropna(how="all")
            .dropna(axis=1, how="all")
        )

        if dataframe.empty:
            continue

        total_rows = len(dataframe)

        for start in range(
            0,
            total_rows,
            rows_per_element,
        ):
            end = min(
                start + rows_per_element,
                total_rows,
            )

            chunk_dataframe = dataframe.iloc[start:end]
            text = dataframe_to_text(chunk_dataframe)

            if not text:
                continue

            elements.append(
                {
                    **metadata,
                    "element_id": (
                        f"{metadata['file_id']}_sheet_"
                        f"{sheet_name}_rows_{start + 1}_{end}"
                    ),
                    "element_type": "xlsx_sheet_rows",
                    "page_number": None,
                    "slide_number": None,
                    "sheet_name": sheet_name,
                    "row_start": start + 1,
                    "row_end": end,
                    "section": sheet_name,
                    "content": text,
                }
            )

    return elements


def parse_file(
    row: dict[str, Any] | pd.Series,
) -> list[dict[str, Any]]:
    """Dispatch one manifest row to the parser for its file type."""
    file_type = safe_text(
        row.get("file_type")
    ).lower()

    file_path = Path(
        safe_text(row.get("absolute_path"))
    )

    if file_type not in SUPPORTED_FILE_TYPES:
        raise ValueError(
            f"Unsupported file type: {file_type or 'missing'}"
        )

    if not file_path.exists():
        raise FileNotFoundError(
            f"Document not found: {file_path}"
        )

    if file_type == "pdf":
        return parse_pdf(file_path, row)

    if file_type == "docx":
        return parse_docx(file_path, row)

    if file_type == "pptx":
        return parse_pptx(file_path, row)

    if file_type in {"md", "txt"}:
        return parse_text_file(file_path, row)

    if file_type == "xlsx":
        return parse_xlsx(file_path, row)

    return []


# -----------------------------------------------------------------------------
# Error recording and scoped replacement helpers
# -----------------------------------------------------------------------------
def build_error_record(
    row: dict[str, Any] | pd.Series,
    error_message: str,
    traceback_text: str = "",
    error_type: str = "parse_error",
) -> dict[str, Any]:
    """Build a consistent parsing error record for CSV output."""
    return {
        "file_id": safe_text(row.get("file_id")),
        "opportunity_id": optional_text(
            row.get("opportunity_id")
        ),
        "project_id": optional_text(
            row.get("project_id")
        ),
        "source_file": safe_text(
            row.get("source_file")
        ),
        "relative_path": safe_text(
            row.get("relative_path")
        ),
        "file_type": safe_text(
            row.get("file_type")
        ),
        "database_target": safe_text(
            row.get("database_target")
        ),
        "error_type": error_type,
        "error_message": error_message,
        "traceback": traceback_text,
    }


def _normalize_manifest_rows(
    manifest_rows: Iterable[dict[str, Any] | pd.Series],
) -> list[dict[str, Any]]:
    """Convert mixed manifest row objects into plain dictionaries."""
    normalized: list[dict[str, Any]] = []

    for row in manifest_rows:
        if isinstance(row, pd.Series):
            normalized.append(row.to_dict())
        else:
            normalized.append(dict(row))

    return normalized


def _replace_opportunity_records(
    existing_rows: list[dict[str, Any]],
    new_rows: list[dict[str, Any]],
    opportunity_id: str,
    database_target: str,
) -> list[dict[str, Any]]:
    """Replace records for one opportunity while preserving all other rows."""
    target_id = safe_text(opportunity_id)
    target_database = safe_text(database_target)

    preserved = [
        row
        for row in existing_rows
        if not (
            safe_text(row.get("opportunity_id"))
            == target_id
            and safe_text(row.get("database_target"))
            == target_database
        )
    ]

    return preserved + new_rows


def _save_errors(
    errors: list[dict[str, Any]],
    errors_path: Path,
    opportunity_id: str | None = None,
    database_target: str | None = None,
    replace_existing: bool = False,
) -> None:
    """Persist parsing errors, optionally replacing one opportunity's errors."""
    errors_path.parent.mkdir(parents=True, exist_ok=True)

    existing_errors: list[dict[str, Any]] = []

    if (
        errors_path.exists()
        and errors_path.stat().st_size > 0
    ):
        try:
            existing_dataframe = pd.read_csv(
                errors_path,
                dtype=str,
                keep_default_na=False,
            )

            existing_errors = (
                existing_dataframe.to_dict(
                    orient="records"
                )
                if not existing_dataframe.empty
                else []
            )

        except pd.errors.EmptyDataError:
            existing_errors = []

    if (
        replace_existing
        and opportunity_id is not None
        and database_target is not None
    ):
        all_errors = _replace_opportunity_records(
            existing_rows=existing_errors,
            new_rows=errors,
            opportunity_id=opportunity_id,
            database_target=database_target,
        )
    else:
        all_errors = existing_errors + errors

    columns = [
        "file_id",
        "opportunity_id",
        "project_id",
        "source_file",
        "relative_path",
        "file_type",
        "database_target",
        "error_type",
        "error_message",
        "traceback",
    ]

    pd.DataFrame(
        all_errors,
        columns=columns,
    ).to_csv(
        errors_path,
        index=False,
        encoding="utf-8-sig",
    )


def parse_opportunity_documents(
    manifest_rows: Iterable[dict[str, Any] | pd.Series],
    opportunity_id: str,
    output_path: Path = OUTPUT_JSONL,
    errors_path: Path = ERRORS_CSV,
    database_target: str = "rfp_db",
    replace_existing: bool = True,
    show_progress: bool = False,
) -> dict[str, Any]:
    """
    Parse only the active opportunity's manifest rows.

    Existing parsed elements for other opportunities and proposal
    knowledge are preserved. When replace_existing=True, records for
    this opportunity are replaced instead of duplicated.
    """
    target_id = safe_text(opportunity_id)

    if not target_id:
        raise ValueError(
            "Opportunity ID cannot be empty."
        )

    rows = _normalize_manifest_rows(manifest_rows)

    selected_rows = [
        row
        for row in rows
        if (
            safe_text(row.get("opportunity_id"))
            == target_id
            and safe_text(row.get("database_target"))
            == database_target
        )
    ]

    if not selected_rows:
        raise ValueError(
            "No manifest rows matched opportunity_id="
            f"'{target_id}' and database_target="
            f"'{database_target}'."
        )

    new_elements: list[dict[str, Any]] = []
    errors: list[dict[str, Any]] = []

    # Streamlit uses the same parser without progress output, while CLI runs can
    # show tqdm progress for longer batches.
    iterator: Iterable[dict[str, Any]] = selected_rows

    if show_progress:
        iterator = tqdm(
            selected_rows,
            total=len(selected_rows),
            desc=f"Parsing {target_id}",
        )

    for row in iterator:
        try:
            elements = parse_file(row)

            if not elements:
                file_type = safe_text(
                    row.get("file_type")
                ).lower()

                if file_type == "pdf":
                    errors.append(
                        build_error_record(
                            row=row,
                            error_type="needs_ocr",
                            error_message=(
                                "PDF contained no extractable text. "
                                "OCR may be required."
                            ),
                        )
                    )
                else:
                    errors.append(
                        build_error_record(
                            row=row,
                            error_type="empty_document",
                            error_message=(
                                "No extractable content was found."
                            ),
                        )
                    )

            new_elements.extend(elements)

        except Exception as error:
            errors.append(
                build_error_record(
                    row=row,
                    error_message=str(error),
                    traceback_text=traceback.format_exc(),
                )
            )

    existing_elements = read_jsonl(output_path)

    if replace_existing:
        all_elements = _replace_opportunity_records(
            existing_rows=existing_elements,
            new_rows=new_elements,
            opportunity_id=target_id,
            database_target=database_target,
        )
    else:
        all_elements = existing_elements + new_elements

    write_jsonl(
        rows=all_elements,
        output_path=output_path,
    )

    _save_errors(
        errors=errors,
        errors_path=errors_path,
        opportunity_id=target_id,
        database_target=database_target,
        replace_existing=replace_existing,
    )

    return {
        "opportunity_id": target_id,
        "database_target": database_target,
        "files_selected": len(selected_rows),
        "parsed_element_count": len(new_elements),
        "error_count": len(errors),
        "elements": new_elements,
        "errors": errors,
        "output_path": str(output_path),
        "errors_path": str(errors_path),
    }


def parse_all_manifest_documents(
    manifest_path: Path = MANIFEST_PATH,
    output_path: Path = OUTPUT_JSONL,
    errors_path: Path = ERRORS_CSV,
) -> dict[str, Any]:
    """Parse every manifest row and rebuild the complete parsed-elements file."""
    if not manifest_path.exists():
        raise FileNotFoundError(
            f"Manifest not found: {manifest_path}. "
            "Run create_manifest.py first."
        )

    manifest = pd.read_csv(manifest_path)

    all_elements: list[dict[str, Any]] = []
    errors: list[dict[str, Any]] = []

    iterator = tqdm(
        manifest.iterrows(),
        total=len(manifest),
        desc="Parsing documents",
    )

    for _, row in iterator:
        try:
            elements = parse_file(row)
            all_elements.extend(elements)

            if not elements:
                file_type = safe_text(
                    row.get("file_type")
                ).lower()

                errors.append(
                    build_error_record(
                        row=row,
                        error_type=(
                            "needs_ocr"
                            if file_type == "pdf"
                            else "empty_document"
                        ),
                        error_message=(
                            "PDF contained no extractable text. "
                            "OCR may be required."
                            if file_type == "pdf"
                            else "No extractable content was found."
                        ),
                    )
                )

        except Exception as error:
            errors.append(
                build_error_record(
                    row=row,
                    error_message=str(error),
                    traceback_text=traceback.format_exc(),
                )
            )

    write_jsonl(
        rows=all_elements,
        output_path=output_path,
    )

    pd.DataFrame(errors).to_csv(
        errors_path,
        index=False,
        encoding="utf-8-sig",
    )

    return {
        "files_selected": len(manifest),
        "parsed_element_count": len(all_elements),
        "error_count": len(errors),
        "elements": all_elements,
        "errors": errors,
        "output_path": str(output_path),
        "errors_path": str(errors_path),
    }


def print_result_summary(result: dict[str, Any]) -> None:
    """Print parser output counts and element-type summaries."""
    print(
        "Parsed elements saved to:",
        result["output_path"],
    )
    print(
        "Parsing errors saved to:",
        result["errors_path"],
    )
    print(
        "Files selected:",
        result["files_selected"],
    )
    print(
        "Total parsed elements:",
        result["parsed_element_count"],
    )
    print(
        "Total errors:",
        result["error_count"],
    )

    elements = result.get("elements", [])

    if elements:
        preview = pd.DataFrame(elements)

        print("\nParsed element types:")
        print(preview["element_type"].value_counts())

        print("\nDatabase targets:")
        print(
            preview["database_target"].value_counts(
                dropna=False
            )
        )

        print("\nCollections:")
        print(
            preview["collection_name"].value_counts(
                dropna=False
            )
        )


def main() -> None:
    """Run the parser from the command line."""
    parser = argparse.ArgumentParser()

    parser.add_argument(
        "--manifest",
        default=str(MANIFEST_PATH),
    )

    parser.add_argument(
        "--output",
        default=str(OUTPUT_JSONL),
    )

    parser.add_argument(
        "--errors",
        default=str(ERRORS_CSV),
    )

    parser.add_argument(
        "--opportunity-id",
        default="",
    )

    parser.add_argument(
        "--database-target",
        default="rfp_db",
    )

    args = parser.parse_args()

    manifest_path = Path(args.manifest)
    output_path = Path(args.output)
    errors_path = Path(args.errors)

    if args.opportunity_id:
        if not manifest_path.exists():
            raise FileNotFoundError(
                f"Manifest not found: {manifest_path}. "
                "Run create_manifest.py first."
            )

        manifest = pd.read_csv(
            manifest_path,
            dtype=str,
            keep_default_na=False,
        )

        result = parse_opportunity_documents(
            manifest_rows=manifest.to_dict(
                orient="records"
            ),
            opportunity_id=args.opportunity_id,
            output_path=output_path,
            errors_path=errors_path,
            database_target=args.database_target,
            replace_existing=True,
            show_progress=True,
        )
    else:
        result = parse_all_manifest_documents(
            manifest_path=manifest_path,
            output_path=output_path,
            errors_path=errors_path,
        )

    print_result_summary(result)


if __name__ == "__main__":
    main()
